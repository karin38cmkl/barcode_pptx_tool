import io
import re
import csv

import fitz
import requests
import streamlit as st
from PIL import Image
from lxml import etree
from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ─── ページ設定 ───────────────────────────────────────────
st.set_page_config(page_title="バーコードPPTX生成ツール", layout="wide")
st.title("バーコードラベル PPTX生成ツール")

# ─── セッションステート初期化 ───────────────────────────────
if "blocks" not in st.session_state:
    st.session_state.blocks = [
        {"type": "変数", "value": "K"},
        {"type": "固定", "value": "_"},
        {"type": "変数", "value": "H"},
        {"type": "固定", "value": "("},
        {"type": "変数", "value": "G"},
        {"type": "固定", "value": ")_"},
        {"type": "変数", "value": "L"},
    ]

# ─── サイドバー設定 ────────────────────────────────────────
st.sidebar.header("設定")

line1_text = st.sidebar.text_input("1行目テキスト（固定）", "MimoRhea(みもれあ) 29×29")

st.sidebar.subheader("2行目テキスト構成")
st.sidebar.caption("変数＝スプシの列、固定＝そのまま出力するテキスト")

# ブロック一覧表示
delete_idx = None
for i, block in enumerate(st.session_state.blocks):
    c1, c2, c3 = st.sidebar.columns([2, 3, 1])
    new_type = c1.selectbox(
        "", ["変数", "固定"],
        index=0 if block["type"] == "変数" else 1,
        key=f"type_{i}",
        label_visibility="collapsed"
    )
    new_val = c2.text_input(
        "", block["value"],
        key=f"val_{i}",
        label_visibility="collapsed",
        placeholder="列名(例:K)" if block["type"] == "変数" else "固定テキスト"
    )
    if c3.button("✕", key=f"del_{i}"):
        delete_idx = i
    st.session_state.blocks[i]["type"] = new_type
    st.session_state.blocks[i]["value"] = new_val

if delete_idx is not None:
    st.session_state.blocks.pop(delete_idx)
    st.rerun()

# 追加ボタン
add1, add2 = st.sidebar.columns(2)
if add1.button("＋ 変数"):
    st.session_state.blocks.append({"type": "変数", "value": ""})
    st.rerun()
if add2.button("＋ 固定"):
    st.session_state.blocks.append({"type": "固定", "value": ""})
    st.rerun()

# プレビュー
preview_parts = []
for b in st.session_state.blocks:
    if b["type"] == "固定":
        preview_parts.append(b["value"])
    else:
        preview_parts.append(f"[{b['value']}列]")
st.sidebar.caption(f"プレビュー: {''.join(preview_parts)}")

st.sidebar.divider()

# ファイル名照合用の商品コード列
code_col = st.sidebar.text_input("商品コード列（ファイル名照合用）", "K",
                                  help="画像ファイル名(A-01など)と照合するための列")
header_rows = st.sidebar.number_input("ヘッダー行数（スキップ行数）", min_value=0, max_value=5, value=1)

st.sidebar.subheader("ページ・デザイン設定")
page_w_mm = st.sidebar.number_input("ページ幅 (mm)", value=400, min_value=50, max_value=1000)
page_h_mm = st.sidebar.number_input("ページ高さ (mm)", value=200, min_value=50, max_value=1000)
white_box_pt = st.sidebar.number_input("白塗り高さ (pt)", value=193, min_value=10, max_value=500)
font_pt = st.sidebar.number_input("フォントサイズ (pt)", value=60, min_value=8, max_value=200)
scale = st.sidebar.slider("画像解像度 (PDF拡大率)", min_value=2, max_value=10, value=7,
                          help="7 = 700%相当。大きいほど高画質だが処理が遅い")

# ─── ヘルパー関数 ──────────────────────────────────────────

def col_letter_to_index(letter: str) -> int:
    letter = letter.strip().upper()
    idx = 0
    for ch in letter:
        idx = idx * 26 + (ord(ch) - ord('A') + 1)
    return idx - 1


def make_filename(code: str) -> str:
    parts = code.split('_')
    if len(parts) >= 2:
        return f"{parts[0][0]}-{parts[1]}"
    return code


def generate_line2(row: list, blocks: list) -> str:
    result = ""
    for block in blocks:
        if block["type"] == "固定":
            result += block["value"]
        else:
            idx = col_letter_to_index(block["value"])
            result += row[idx].strip() if len(row) > idx else ""
    return result


def remove_shadow(shape):
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        spPr = sp.find('spPr')
    if spPr is not None:
        for el in spPr.findall(qn('a:effectLst')):
            spPr.remove(el)
        etree.SubElement(spPr, qn('a:effectLst'))


def add_textbox(slide, text, left, top, width, height, font_pt, bold=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.bold = bold
    r.font.size = Pt(font_pt)
    r.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    remove_shadow(txBox)
    return txBox


def extract_barcodes_from_pdf(pdf_bytes: bytes, scale: int) -> dict:
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    results = {}

    for page in doc:
        mat = fitz.Matrix(scale, scale)
        pix = page.get_pixmap(matrix=mat)
        page_img = Image.open(io.BytesIO(pix.tobytes('png')))
        blocks = page.get_text('blocks')

        product_blocks, barcode_num_blocks = [], []
        for b in blocks:
            text = b[4].strip()
            first_line = text.split('\n')[0].strip()
            if re.match(r'^[A-Z]\d+_\d+_\d+x\d+', first_line):
                product_blocks.append(b)
            elif re.match(r'^\d{10,15}$', text.replace('\n', '')):
                barcode_num_blocks.append(b)

        for pb in product_blocks:
            first_line = pb[4].strip().split('\n')[0].strip()
            filename = make_filename(first_line)
            px0, py0, px2, py2 = pb[0], pb[1], pb[2], pb[3]
            pc = (px0 + px2) / 2

            matched_nb, min_dist = None, float('inf')
            for nb in barcode_num_blocks:
                nc = (nb[0] + nb[2]) / 2
                if nb[1] > py0 and abs(nc - pc) < 30:
                    dist = nb[1] - py2
                    if dist < min_dist:
                        min_dist = dist
                        matched_nb = nb

            if matched_nb is None:
                continue

            cropped = page_img.crop((
                int(min(px0, matched_nb[0]) - 2) * scale,
                int(py0 - 2) * scale,
                int(max(px2, matched_nb[2]) + 2) * scale,
                int(matched_nb[3] + 2) * scale,
            ))
            results[filename] = cropped

    doc.close()
    return results


def load_spreadsheet_data(url, csv_bytes, code_col_idx, blocks, header_rows) -> dict:
    if url:
        match = re.search(r'/spreadsheets/d/([a-zA-Z0-9_-]+)', url)
        if not match:
            st.error("Google SheetsのURLが正しくありません")
            return {}
        sheet_id = match.group(1)
        export_url = f"https://docs.google.com/spreadsheets/d/{sheet_id}/export?format=csv"
        resp = requests.get(export_url, allow_redirects=True, timeout=30)
        rows = list(csv.reader(io.StringIO(resp.content.decode('utf-8'))))
    elif csv_bytes:
        rows = list(csv.reader(io.StringIO(csv_bytes.decode('utf-8'))))
    else:
        return {}

    data = {}
    for row in rows[header_rows:]:
        code = row[code_col_idx].strip() if len(row) > code_col_idx else ''
        if code and code != '-':
            data[code] = generate_line2(row, blocks)
    return data


def build_pptx(barcode_images, ss_data, line1, page_w_mm, page_h_mm,
               white_box_pt, font_pt) -> bytes:
    prs = Presentation()
    prs.slide_width = Mm(page_w_mm)
    prs.slide_height = Mm(page_h_mm)
    blank = prs.slide_layouts[6]

    white_h = Pt(white_box_pt)
    half_h = white_h // 2

    for code in sorted(barcode_images.keys()):
        if code not in ss_data:
            continue

        img = barcode_images[code]
        w_px, h_px = img.size
        img_h_mm = page_w_mm * h_px / w_px

        slide = prs.slides.add_slide(blank)

        buf = io.BytesIO()
        img.save(buf, format='JPEG', quality=95)
        buf.seek(0)
        slide.shapes.add_picture(buf, Mm(0), Mm(0), Mm(page_w_mm), Mm(img_h_mm))

        rect = slide.shapes.add_shape(1, Mm(0), Mm(0), Mm(page_w_mm), white_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        rect.line.fill.background()
        remove_shadow(rect)

        add_textbox(slide, line1,        Mm(0), Mm(0),   Mm(page_w_mm), half_h, font_pt)
        add_textbox(slide, ss_data[code], Mm(0), half_h, Mm(page_w_mm), half_h, font_pt)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ─── メイン UI ─────────────────────────────────────────────

col1, col2 = st.columns(2)

with col1:
    st.subheader("① PDFをアップロード")
    pdf_file = st.file_uploader("バーコードPDF", type=["pdf"])

with col2:
    st.subheader("② スプレッドシートデータ")
    tab_url, tab_csv = st.tabs(["Google Sheets URL", "CSVアップロード"])
    with tab_url:
        sheets_url = st.text_input("URLを貼り付け",
                                   placeholder="https://docs.google.com/spreadsheets/d/...")
    with tab_csv:
        csv_file = st.file_uploader("CSVファイル", type=["csv"])

st.divider()

if st.button("▶ PPTXを生成", type="primary", use_container_width=True):
    if not pdf_file:
        st.error("PDFをアップロードしてください")
        st.stop()

    use_url = sheets_url.strip() if sheets_url else None
    use_csv = csv_file.read() if csv_file else None

    if not use_url and not use_csv:
        st.error("スプレッドシートのURLまたはCSVを指定してください")
        st.stop()

    code_col_idx = col_letter_to_index(code_col)
    blocks_snapshot = [dict(b) for b in st.session_state.blocks]

    with st.spinner("PDFからバーコード画像を切り出し中..."):
        try:
            barcode_images = extract_barcodes_from_pdf(pdf_file.read(), scale)
            st.success(f"バーコード画像: {len(barcode_images)}件 切り出し完了")
        except Exception as e:
            st.error(f"PDF処理エラー: {e}")
            st.stop()

    with st.spinner("スプレッドシートデータを読み込み中..."):
        try:
            ss_data = load_spreadsheet_data(use_url, use_csv, code_col_idx,
                                            blocks_snapshot, int(header_rows))
            st.success(f"スプレッドシート: {len(ss_data)}件 読み込み完了")
        except Exception as e:
            st.error(f"スプレッドシート読み込みエラー: {e}")
            st.stop()

    matched  = set(barcode_images.keys()) & set(ss_data.keys())
    only_img = set(barcode_images.keys()) - set(ss_data.keys())
    only_ss  = set(ss_data.keys()) - set(barcode_images.keys())

    if only_img:
        st.warning(f"スプシにデータなし（スキップ）: {sorted(only_img)}")
    if only_ss:
        st.warning(f"PDF画像なし（スキップ）: {sorted(only_ss)}")

    with st.spinner(f"PPTX生成中（{len(matched)}枚）..."):
        try:
            pptx_bytes = build_pptx(
                barcode_images, ss_data, line1_text,
                int(page_w_mm), int(page_h_mm),
                int(white_box_pt), int(font_pt)
            )
            st.success(f"PPTX生成完了: {len(matched)}スライド")
        except Exception as e:
            st.error(f"PPTX生成エラー: {e}")
            st.stop()

    st.download_button(
        label="⬇ PPTXをダウンロード",
        data=pptx_bytes,
        file_name="barcode_labels.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        use_container_width=True,
        type="primary"
    )

    st.subheader("プレビュー（最初の3件）")
    preview_codes = sorted(matched)[:3]
    cols = st.columns(3)
    for i, code in enumerate(preview_codes):
        with cols[i]:
            st.image(barcode_images[code], caption=code, use_container_width=True)
            st.caption(ss_data[code])
